"""
Tools to make contexts (knowledge bases for AI agents)

Examples:

Download all articles from a markdown string and save them as PDF files:

>>> download_articles(md_string)  # doctest: +SKIP

Verify URLs in a markdown string by checking their status codes
(useful when trying to verify if AI hallucinated the urls)

>>> verify_urls(md_string)  # doctest: +SKIP

Make an md file with all the code in a directory:

>>> md_string = code_aggregate(package_or_folder_or_github_url)  # doctest: +SKIP

"""

# --------------------------------------------------------------------------------------
# General utilities

import os
from typing import Union, Callable, Optional, Iterable
from types import ModuleType
from functools import partial
from pathlib import Path
from dol import written_key


def identity(x):
    return x


def fullpath(path: str) -> str:
    """
    Returns the full path of the given path.

    Args:
        path (str): The path to convert to a full path.

    Returns:
        str: The full path.

    Example:

    >>> fullpath('~/Downloads')  # doctest: +SKIP
    '/home/user/Downloads'

    >>> fullpath('.')  # doctest: +SKIP
    '/home/user/python_projects/aix/aix'

    """
    return os.path.abspath(os.path.expanduser(path))


def save_to_file_and_return_file(
    obj=None, *, encoder=identity, key: Union[str, Callable] = None
):
    """
    Save `encoder(obj)` to a file using a random name in `rootdir` (or a temp directory if not provided).
    Returns the full path to the saved file.
    If `obj` is None, returns a partial function with preconfigured `encoder` and
    `rootdir`.

    :param obj: The object to save. If None, return a partial function.
    :param encoder: A function to encode the object into text or bytes. Defaults to identity.
    :param key: The key (by default, filepath) to write to.
        If None, a temporary file is created.
        If a string starting with '*', the '*' is replaced with a unique temporary filename.
        If a string that has a '*' somewhere in the middle, what's on the left of if is used as a directory
        and the '*' is replaced with a unique temporary filename. For example
        '/tmp/*_file.ext' would be replaced with '/tmp/oiu8fj9873_file.ext'.
        If a callable, it will be called with obj as input to get the key. One use case
        is to use a function that generates a key based on the object.

    :return: Full path to the saved file, or a partial function if `obj` is None.

    >>> from pathlib import Path
    >>> filepath = save_to_file_and_return_file("hello world")
    >>> import os
    >>> Path(filepath).read_text()
    'hello world'

    The default encoder is identity, so you can save binary data as well:

    >>> filepath = save_to_file_and_return_file(b"binary data", encoder=lambda x: x)
    >>> Path(filepath).read_bytes()
    b'binary data'

    But when your object is neither text nor bytes, you can specify a custom encoder
    that transforms your object into text or bytes.
    Also, see how if you don't specify the obj to save, you get a partial function that
    you can use later. This is useful when you want to save many objects with the same
    encoder.

    See below how to make a json_save_and_get_path

    >>> import json
    >>> json_save_and_get_path = save_to_file_and_return_file(encoder=json.dumps)
    >>> filepath = json_save_and_get_path({"key": "value"})
    >>> json.loads(Path(filepath).read_text())
    {'key': 'value'}

    """
    # Note: Yes, it's just written_key from dol, but with a context-sensitive name
    return written_key(obj, encoder=encoder, key=key)


# --------------------------------------------------------------------------------------
# Converting things to markdown

import contextlib
import base64
import io
import json
from typing import Callable, Dict, Optional, Mapping, MutableMapping, Any

from dol import Files, TextFiles

# Import libraries with error suppression
ignore_import_errors = contextlib.suppress(ImportError)

# Initialize the default converters dictionary
dflt_converters: Dict[str, Callable[[bytes], str]] = {}

dflt_md_inner_file_header = "###"

# PDF Conversion
with ignore_import_errors:
    import pypdf

    def pdf_to_markdown(
        pdf_bytes: bytes, *, md_inner_file_header=dflt_md_inner_file_header
    ) -> str:
        """Convert PDF to markdown text."""
        pdf_reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
        pages = []
        for page in pdf_reader.pages:
            text = page.extract_text()
            pages.append(f"{md_inner_file_header} Page {len(pages) + 1}\n\n{text}")
        return "\n\n".join(pages)

    dflt_converters["pdf"] = pdf_to_markdown

# Microsoft Office Conversion
with ignore_import_errors:
    import mammoth  # for .docx

    def docx_to_markdown(docx_bytes: bytes) -> str:
        """Convert DOCX to markdown."""
        result = mammoth.convert_to_markdown(io.BytesIO(docx_bytes))
        return result.value

    dflt_converters["docx"] = docx_to_markdown
    dflt_converters["doc"] = docx_to_markdown

# Excel Conversion
with ignore_import_errors:
    import pandas as pd

    def excel_to_markdown(
        excel_bytes: bytes, *, md_inner_file_header=dflt_md_inner_file_header
    ) -> str:
        """Convert Excel files to markdown tables."""
        # Support both xls and xlsx
        dfs = pd.read_excel(io.BytesIO(excel_bytes), sheet_name=None)
        markdown_output = []

        for sheet_name, df in dfs.items():
            markdown_output.append(f"{md_inner_file_header} Sheet: {sheet_name}\n")
            # Convert DataFrame to markdown table
            markdown_output.append(df.to_markdown(index=False))
            markdown_output.append("\n")

        return "\n".join(markdown_output)

    dflt_converters["xlsx"] = excel_to_markdown
    dflt_converters["xls"] = excel_to_markdown

# PowerPoint Conversion
with ignore_import_errors:
    import pptx

    def pptx_to_markdown(
        pptx_bytes: bytes, *, md_inner_file_header=dflt_md_inner_file_header
    ) -> str:
        """Convert PowerPoint to markdown."""
        presentation = pptx.Presentation(io.BytesIO(pptx_bytes))
        slides = []

        for i, slide in enumerate(presentation.slides, 1):
            slide_text = [f"{md_inner_file_header} Slide {i}\n"]

            # Extract text from various elements
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            slides.append("\n".join(slide_text))

        return "\n\n".join(slides)

    dflt_converters["pptx"] = pptx_to_markdown
    dflt_converters["ppt"] = pptx_to_markdown

# HTML Conversion
with ignore_import_errors:
    import html2text

    def html_to_markdown(html_bytes: bytes) -> str:
        """Convert HTML to markdown."""
        # Decode bytes to string
        html_str = html_bytes.decode("utf-8", errors="ignore")

        # Create HTML to Markdown converter
        h = html2text.HTML2Text()
        h.ignore_links = False
        h.ignore_images = False

        return h.handle(html_str)

    dflt_converters["html"] = html_to_markdown


# Fallback converter (default)
def default_fallback(b: bytes, input_format: str) -> str:
    """Fallback converter that attempts basic text extraction."""
    try:
        # Try decoding with different encodings
        encodings = ["utf-8", "latin-1", "ascii", "utf-16"]
        for encoding in encodings:
            try:
                return b.decode(encoding)
            except UnicodeDecodeError:
                continue

        # If decoding fails, convert to base64 for preservation
        return f"# Conversion Failed for {input_format}\n\n```\n{base64.b64encode(b).decode('utf-8')}\n```"
    except Exception as e:
        return f"# Conversion Error\n\nCould not convert {input_format}: {str(e)}"


def convert_to_markdown(
    b: bytes,
    input_format: str,
    *,
    converters: Dict[str, Callable[[bytes], str]] = dflt_converters,
    fallback: Callable = default_fallback,
) -> str:
    """
    Convert bytes of a given format to markdown text.

    Args:
        b (bytes): Input bytes to convert
        input_format (str): Format of the input (e.g., 'pdf', 'docx')
        converters (Dict[str, Callable]): Dictionary of format-specific converters
        fallback (Callable): Fallback conversion method

    Returns:
        str: Markdown-formatted text
    """
    converter = converters.get(input_format.lower(), None)
    if converter is not None:
        return converter(b)
    else:
        return fallback(b, input_format)


def add_dflt_converter(input_format: str, converter: Callable[[bytes], str]):
    """Add (or change) a default converter for a given input format."""
    dflt_converters[input_format.lower()] = converter


def get_extension(path: str) -> str:
    """Get the extension of a file path."""
    return Path(path).suffix.lstrip(".").lower()


def extensions_not_supported_by_converters(
    paths: Union[str, Iterable[str]], converters: Iterable = dflt_converters
):
    """
    Returns a set of file extensions that are not supported by the given converters.

    Args:
        paths (str or Iterable): A list of file paths or the root directory to check for paths.
        converters (dict): An iterable of extensions supported by the converters.
            By default, it's the default converters (a dict, so keys are considered).

    Returns:
        set: A set of file extensions that are not supported by the given converters.

    """
    if isinstance(paths, str):
        paths = Files(paths)

    paths_extensions = {get_extension(path) for path in paths}
    supported_extensions = set(converters)

    return paths_extensions - supported_extensions


def _resolve_src_bytes_store_and_target_text_store(src_files, target_store):
    """
    Resolves the source bytes store and target text store.
    """
    if isinstance(src_files, str):
        src_rootdir = fullpath(src_files)
        src_files = Files(src_rootdir)
    else:
        if hasattr(src_files, "rootdir"):
            src_rootdir = src_files.rootdir
        else:
            src_rootdir = None

    if target_store is None:
        target_store = dict()

    if isinstance(target_store, str):
        target_store = TextFiles(target_store)

    return src_files, target_store


def convert_files_to_markdown(
    src_files: Union[str, Mapping[str, bytes]],
    target_store: Union[str, MutableMapping[str, bytes]] = None,
    *,
    converters: Dict[str, Callable[[bytes], str]] = dflt_converters,
    verbose: bool = False,
    old_to_new_key: Callable[[str], str] = lambda x: x + ".md",
    target_store_egress: Callable[[Mapping], Any] = identity,
):
    """
    Converts files to markdown using the specified converters.

    Tip: You can use `target_store_egress=aggregate_store` to get an aggregate string
    as your output.
    """
    src_files, target_store = _resolve_src_bytes_store_and_target_text_store(
        src_files, target_store
    )
    _convert_to_markdown = partial(convert_to_markdown, converters=converters)

    for src_key, src_bytes in src_files.items():
        if verbose:
            print(f"Converting {src_key} to markdown")
        ext = get_extension(src_key)
        target_key = old_to_new_key(src_key)
        target_store[target_key] = _convert_to_markdown(src_bytes, ext)

    return target_store_egress(target_store)


# --------------------------------------------------------------------------------------
# Store aggregation

from typing import Optional, Mapping, Iterable, Union
from dol import TextFiles, filt_iter, wrap_kvs, Pipe, store_aggregate


def aggregate_store(
    store: Mapping,
    *,
    min_number_of_duplicated_lines: Optional[int] = None,
    max_num_characters: Optional[int] = None,
    exclude: Optional[Iterable] = None,
    chk_size=None,
    egress: Optional[Union[str, Callable]] = None,
    **store_aggregate_kwargs,
):
    """
    Aggregates and processes text files from a store.

    Parameters:
        store: The source store (e.g., TextFiles instance)
        min_number_of_duplicated_lines (int): If not None, minimum block size for deduplication.
        max_num_characters (int): Maximum number of characters per file.
        exclude (set): Set of filenames to exclude.
        chk_size (int): Chunk size for aggregation. If None, won't chunk
        egress (str or callable): Template for output filenames or function to call on each chunk
        **store_aggregate_kwargs: Additional keyword arguments to pass to store_aggregate.
    """
    if exclude is None:
        exclude = set()

    wrappers = []

    if exclude is not None:
        wrappers.append(filt_iter(filt=lambda x: x not in exclude))

    if min_number_of_duplicated_lines is not None:
        from hg import deduplicate_string_lines

        remove_duplicate_lines = wrap_kvs(
            value_decoder=lambda v: deduplicate_string_lines(
                v,
                min_block_size=min_number_of_duplicated_lines,
                return_removed_blocks=False,
            )
        )
    else:
        remove_duplicate_lines = identity

    if max_num_characters is not None:
        wrappers.append(wrap_kvs(value_decoder=lambda v: v[:max_num_characters]))

    capped_num_characters = wrap_kvs(value_decoder=lambda v: v[:max_num_characters])

    if len(wrappers) > 0:
        store_wrap = Pipe(
            filt_iter(filt=lambda x: x.endswith(".md") and x not in exclude),
            remove_duplicate_lines,
            capped_num_characters,
        )
    else:
        store_wrap = identity

    wrapped_store = store_wrap(store)

    if isinstance(egress, str):
        output_template = egress
        egress = output_template.format
    else:
        assert egress is None or callable(
            egress
        ), "egress must be None, str, or callable"

    if isinstance(chk_size, int):
        from lkj.chunking import chunk_iterable

        chunks = chunk_iterable(wrapped_store, chk_size)
        if egress is not None:
            egress = "store_aggregate_{:02.0f}.md"
        for i, sub_store in enumerate(chunks, 1):
            store_aggregate(sub_store, egress=egress(i), **store_aggregate_kwargs)
    else:
        if egress is None:
            egress = identity
        return store_aggregate(wrapped_store, egress=egress, **store_aggregate_kwargs)


# --------------------------------------------------------------------------------------
# Convert Jupyter notebooks to Markdown

import os
from typing import Sequence, Callable
from contextlib import suppress

from dol import Pipe
import lkj

ignore_import_errors = suppress(ImportError, ModuleNotFoundError)

installed_packages = set()

with ignore_import_errors:
    from nbconvert.preprocessors import Preprocessor
    from nbconvert import MarkdownExporter

    installed_packages.add("nbconvert")

with ignore_import_errors:
    import nbformat

    installed_packages.add("nbformat")


def truncate_text(text: str, *, max_chars=200) -> str:
    if len(text) > max_chars:
        cutoff = max_chars
        first_part = text[: int(0.75 * cutoff)]
        last_part = text[-int(0.25 * cutoff) :]
        return f"{first_part}\n...\n{last_part}"
    return text


# TODO: Generalize to get more control over output (and source) transformation control
def notebook_to_markdown(
    src_notebook: str,
    *,
    target_file: str = "*.md",
    output_processors: Sequence[Callable[[str], str]] = (truncate_text,),
    read_encoding: str = "utf-8",
    write_encoding: str = "utf-8",
):
    """Converts a Jupyter notebook to Markdown, applying a output text processors."""
    if "nbconvert" not in installed_packages or "nbformat" not in installed_packages:
        raise ImportError("nbconvert and nbformat must both be installed.")

    class TruncateOutputPreprocessor(Preprocessor):
        """Preprocessor that truncates long outputs in Jupyter Notebook cells."""

        def __init__(self, output_processors, **kwargs):
            super().__init__(**kwargs)
            self.output_processors = output_processors
            if self.output_processors:
                self.process_output = Pipe(*output_processors)
            else:
                self.process_output = identity

        def preprocess_cell(self, cell, resources, index):
            if "outputs" in cell:
                for output in cell.outputs:
                    # Case 1: Traditional text output
                    if isinstance(output, dict):
                        if "text" in output:
                            output["text"] = self.process_output(output["text"])

                        # Case 2: Rich MIME-based data (e.g. text/html, text/plain, etc.)
                        if "data" in output:
                            # Clear all the existing fields
                            output.clear()

                            # Convert it to a 'stream'-type output
                            # 'text' is valid only when 'output_type' is 'stream' or 'error'
                            output["output_type"] = "stream"
                            output["name"] = "stdout"  # required for 'stream'
                            output["text"] = "HTML output truncated. (Data removed)\n"
            return cell, resources

    if target_file.startswith("*"):
        src_notebook_no_ext = os.path.splitext(src_notebook)[0]
        target_file = target_file.replace("*", src_notebook_no_ext, 1)

    # Load notebook
    with open(src_notebook, "r", encoding=read_encoding) as f:
        nb = nbformat.read(f, as_version=4)

    # Configure exporter with preprocessor
    exporter = MarkdownExporter()

    exporter.register_preprocessor(
        TruncateOutputPreprocessor(output_processors), enabled=True
    )

    # Convert notebook
    body, _ = exporter.from_notebook_node(nb)

    # Save markdown file
    with open(target_file, "w", encoding=write_encoding) as f:
        f.write(body)

    return target_file


# --------------------------------------------------------------------------------------
# Download articles from a markdown string and save them as PDF files

# TODO: Make download_articles more general:
#       - Allowed file types should be handled by plugin dependency injection
#       - There should be a separate title/url extractor that can be passed in
#       - When title is missing, some url_to_filenae should be used (see graze)?
#       - download_articles_by_section should be merged with download_articles

import os
import re
import requests

DFLT_SAVE_DIR = os.path.expanduser("~/Downloads")


def download_articles(
    md_string: str,
    save_dir: str = DFLT_SAVE_DIR,
    *,
    save_non_pdf: bool = False,
    verbose: bool = True,
):
    """
    Downloads articles from the given markdown string and saves them as PDF files.

    Args:
        md_string (str): The markdown-style string containing titles and URLs.
        save_dir (str): The root directory to save the downloaded PDFs. Defaults to '~/Downloads'.
        save_non_pdf (bool): Whether to save non-PDF content. Defaults to False.
        verbose (bool): Whether to print detailed messages. Defaults to True.

    Returns:
        list: A list of URLs that failed to download or were invalid PDFs.

    Example:

    >>> md_string = '''
    ... - **[Valid PDF](https://example.com/file.pdf)**: A valid PDF file.
    ... - **[Invalid PDF](https://example.com/file.html)**: An HTML page, not a PDF.
    ... '''  # doctest: +SKIP
    >>> download_articles(md_string, save_non_pdf=True)  # doctest: +SKIP
    Downloaded: Valid PDF -> ~/Downloads/Valid_PDF.pdf
    Skipped (HTML or non-PDF): Invalid PDF from https://example.com/file.html
    Non-PDF content saved to: ~/Downloads/Invalid_PDF_non_pdf.html

    Tips:

    - When you knowledge base will have a lot of files, some AI systems have a hard time
        processing the large number of files. In such cases, it might be better to
        aggregate many files into a single file. See pdfdol.concat_pdfs to do this.


    """
    # Assert the save_dir exists
    assert os.path.exists(save_dir), f"Directory not found: {save_dir}"

    def clog(msg):
        if verbose:
            print(msg)

    # Regex to extract titles and URLs from the markdown string
    pattern = r"- \*\*\[(.*?)\]\((.*?)\)\*\*"
    matches = re.findall(pattern, md_string)

    failed_urls = []

    for title, url in matches:
        # Sanitize title to create a valid filename
        sanitized_title = re.sub(r"[^\w\-_\. ]", "_", title)
        filename = f"{sanitized_title}.pdf"
        filepath = os.path.join(save_dir, filename)

        try:
            response = requests.get(url, stream=True)
            response.raise_for_status()

            # Check Content-Type header
            content_type = response.headers.get("Content-Type", "")
            if "application/pdf" not in content_type:
                clog(
                    f"Skipped (HTML or non-PDF): {title} from {url} (Content-Type: {content_type})"
                )
                if save_non_pdf:
                    # Save non-PDF content with a different extension
                    non_pdf_path = os.path.join(
                        save_dir, f"{sanitized_title}_non_pdf.html"
                    )
                    with open(non_pdf_path, "wb") as f:
                        for chunk in response.iter_content(chunk_size=8192):
                            f.write(chunk)
                    clog(f"Non-PDF content saved to: {non_pdf_path}")
                failed_urls.append(url)
                continue

            # Verify PDF content by checking the first few bytes
            first_chunk = next(response.iter_content(chunk_size=8192))
            if not first_chunk.startswith(b"%PDF"):
                clog(f"Invalid PDF content: {title} from {url}")
                if save_non_pdf:
                    # Save invalid PDF content with a different extension
                    invalid_pdf_path = os.path.join(
                        save_dir, f"{sanitized_title}_invalid.pdf"
                    )
                    with open(invalid_pdf_path, "wb") as f:
                        f.write(first_chunk)
                        for chunk in response.iter_content(chunk_size=8192):
                            f.write(chunk)
                    clog(f"Invalid PDF content saved to: {invalid_pdf_path}")
                failed_urls.append(url)
                continue

            # Save the content as a PDF file
            with open(filepath, "wb") as f:
                f.write(first_chunk)  # Write the first chunk already read
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)

            clog(f"Downloaded: {title} -> {filepath}")
        except Exception as e:
            clog(f"Failed to download {title} from {url}: {e}")
            failed_urls.append(url)

    return failed_urls


def download_articles_by_section(
    md_string, rootdir=None, save_non_pdf: bool = False, *, section_marker: str = r"###"
):
    """
    Downloads articles from a markdown string organized by sections into subdirectories.

    This is useful, for example, when you have a large knowledge base and you want to
    organize of aggregate the articles by sections.

    Args:
        md_string (str): The markdown string with sections and articles.
        rootdir (str): The root directory where subdirectories for sections will be created.
                       Defaults to '~/Downloads'.
        save_non_pdf (bool): Whether to save non-PDF content. Defaults to False.

    Returns:
        dict: A dictionary with section names as keys and lists of failed URLs as values.
    """
    if rootdir is None:
        rootdir = os.path.expanduser("~/Downloads")

    # Ensure the root directory exists
    os.makedirs(rootdir, exist_ok=True)

    # Parse sections and their content
    section_pattern = section_marker + r" (.*?)\n(.*?)(?=\n" + section_marker + "|\Z)"
    sections = re.findall(section_pattern, md_string, re.DOTALL)

    failed_urls_by_section = {}

    for section_title, section_content in sections:
        # Create a snake-case directory name for the section
        sanitized_section_title = (
            re.sub(r"[^\w\s]", "", section_title).strip().replace(" ", "_").lower()
        )
        section_dir = os.path.join(rootdir, sanitized_section_title)
        os.makedirs(section_dir, exist_ok=True)

        print(f"\nProcessing section: {section_title} (Directory: {section_dir})")

        # Download articles for this section
        failed_urls = download_articles(
            section_content, save_dir=section_dir, save_non_pdf=save_non_pdf
        )
        failed_urls_by_section[section_title] = failed_urls

    return failed_urls_by_section


def verify_urls(md_string):
    """
    Verifies URLs in a markdown string by checking their status codes.

    Args:
        md_string (str): The markdown string containing URLs.

    Returns:
        dict: A dictionary with URLs as keys and their status codes as values.
    """
    # Regex to extract URLs from the markdown string
    pattern = r"\[(.*?)\]\((.*?)\)"
    matches = re.findall(pattern, md_string)

    url_status_codes = {}

    for title, url in matches:
        try:
            response = requests.head(url, allow_redirects=True)
            url_status_codes[url] = response.status_code
        except Exception as e:
            url_status_codes[url] = str(e)

    return url_status_codes


# --------------------------------------------------------------------------------------
# Code

from typing import Mapping
import os
import importlib
from typing import Any, Union, Callable
from dol import store_aggregate, TextFiles, filt_iter, cached_keys

DirectoryPathString = str
GithubUrl = str
RegexString = str
CodeSource = Union[DirectoryPathString, GithubUrl, Mapping]


def is_local_pkg_name(name: str) -> bool:
    """Returns True if and only if name is the name of a local package."""
    try:
        importlib.import_module(name)
        return True
    except ImportError:
        return False


def resolve_code_source_dir_path(code_src: CodeSource) -> DirectoryPathString:
    """
    Resolves code_src to a directory path string.

    I
    Args:
        code_src (Any): The source of the code. Can be
            a directory path string,
            a GitHub URL,
            or an imported package (must contain a __path__ atribute)

    Returns:
        The resolved directory path.

    Raises:
        AssertionError: If the resolved path is not a valid directory.
    """
    if isinstance(code_src, str):
        # If it's a string, check if it's a directory
        if os.path.isdir(code_src):
            return os.path.abspath(code_src)  # Return absolute path if it's a directory
        elif "\n" not in code_src and "github" in code_src:
            from hubcap import ensure_repo_folder  # pip install hubcap

            # If it's a GitHub URL, download the repository
            repo_url = code_src
            repo_path = ensure_repo_folder(repo_url)
            return repo_path
        elif is_local_pkg_name(code_src):
            code_src = importlib.import_module(code_src)  # import the package
        else:
            raise ValueError(f"Unsupported string format or non-directory: {code_src}")
    # If it's not a string, check for __path__ attribute
    if hasattr(code_src, "__path__"):
        path_list = list(code_src.__path__)
        assert len(path_list) == 1, (
            f"The __path__ attribute should contain exactly one path, "
            f"but found: {path_list}"
        )
        return os.path.abspath(path_list[0])

    # If no valid resolution was found
    raise ValueError(
        f"Unable to resolve code_src to a valid directory path: {code_src}"
    )


def resolve_code_source(
    code_src: CodeSource, keys_filt: Callable = lambda x: x.endswith(".py")
) -> Mapping:
    """
    Will resolve code_src to a Mapping whose values are the code strings

    Args:
        code_src: The source of the code. Can be an explicit `Mapping`,
            a directory path string,
            a GitHub URL,
            or an imported package (must contain a __path__ atribute)
        keys_filt (Callable): A function to filter the keys. Defaults to lambda x: x.endswith('.py').

    """
    if isinstance(code_src, Mapping):
        return code_src
    else:
        code_src_rootdir = resolve_code_source_dir_path(code_src)
        return cached_keys(
            filt_iter(TextFiles(code_src_rootdir), filt=keys_filt), keys_cache=sorted
        )


def identity(x):
    return x


Filepath = str


def code_aggregate(
    code_src: CodeSource,
    *,
    egress: Union[Callable, Filepath] = identity,
    kv_to_item=lambda k, v: f"## {k}\n\n```python\n{v.strip()}\n```",
    **store_aggregate_kwargs,
) -> Any:
    """
    Aggregates all code segments from the given code source (folder, github url, store).

    This is useful when you want to use AI to search and respond to questions about a
    specific code base.

    Args:
        code_src (dict): A dictionary where keys are references to the code (e.g., paths)
                         and values are code snippets or content.
        egress (Union[Callable, str]): A function to apply to the aggregate before returning.
                                       If a string, the aggregate will be saved to the file.
        kv_to_item (Callable): A function that converts a key-value pairs to the
                               items that should be aggregated.
        **store_aggregate_kwargs: Additional keyword arguments to pass to store_aggregate.

    See dol.store_aggregate for more details.

    Returns:
        Any: The aggregated code content, or the result of the egress function.

    Example:

    >>> code_src = {
    ...     'module1.py': 'def foo(): pass',
    ...     'module2.py': 'def bar(): pass',
    ...     'module3.py': 'class Baz: pass',
    ... }
    >>> print(code_aggregate(code_src))
    ## module1.py
    <BLANKLINE>
    ```python
    def foo(): pass
    ```
    <BLANKLINE>
    ## module2.py
    <BLANKLINE>
    ```python
    def bar(): pass
    ```
    <BLANKLINE>
    ## module3.py
    <BLANKLINE>
    ```python
    class Baz: pass
    ```

    Here, let's input an imported (third party) package, and have the function save the
    result to a temporary file

    >>> from tempfile import NamedTemporaryFile
    >>> temp_file_name = NamedTemporaryFile().name
    >>> import aix
    >>> _  = code_aggregate(aix, egress=temp_file_name)
    >>> print(open(temp_file_name).read(15))
    ## __init__.py
    <BLANKLINE>

    Tip: You can also import directly from the name (string) of the package by doing
    `code_aggregate(__import__('aix'))` or more robustly,
    `importlib.import_module('aix')`.

    If you have hubcap installed, you can even get an aggregate of code from a GitHub
    repository.

    >>> string = code_aggregate('https://github.com/thorwhalen/aix')  # doctest: +SKIP


    """
    code_store = resolve_code_source(code_src)
    return store_aggregate(
        code_store, egress=egress, kv_to_item=kv_to_item, **store_aggregate_kwargs
    )


class PackageCodeContexts:
    """Manages aggregation and saves of the code of local packages"""

    def __init__(self, save_folder="."):
        self.save_folder = fullpath(save_folder)
        self.save_filepath = lambda *parts: os.path.join(self.save_folder, *parts)

    def save_single(self, pkg: Union[str, ModuleType]):
        """
        Aggregates and saves the code of a single local package.

        Example:

        To save a single package's code in a single file, in the current folder:

        >>> PackageCodeContexts().save_single('aix')  # doctest: +SKIP

        or, to save multiple package's code in a single file, in a specific folder.

        >>> PackageCodeContexts('some/folder/path').save_multiple_pkgs_code(['aix', 'dol'])  # doctest: +SKIP

        """
        if isinstance(pkg, str):
            pkg_name = pkg
            pkg = importlib.import_module(pkg_name)
        elif isinstance(pkg, ModuleType):
            pkg_name = pkg.__name__
        else:
            raise ValueError(f"Unsupported type for pkg: {pkg}")

        filepath = self.save_filepath(f"{pkg_name}.py.md")
        code_aggregate(pkg, egress=filepath)

    def multiple_pkgs_code(
        self,
        name: Optional[Union[str, Iterable]] = None,
        pkgs: list = (),
        *,
        pkg_secion_marker="#",
        section_sepator="\n\n\n",
    ):
        """
        Aggregates the code of multiple local packages in separate sections.
        Saves the result to a file if a name is provided.
        """
        if (
            name is not None
            and not isinstance(name, str)
            and isinstance(name, Iterable)
            and not pkgs
        ):
            # if the first argument is an iterable, assume it's the list of packages
            pkgs = name  # the first argument is actually packages
            name = None

        def sections():
            for pkg in pkgs:
                yield f"{pkg_secion_marker} {pkg}\n\n" + code_aggregate(pkg)

        md_string = section_sepator.join(sections())

        if name is None:
            return md_string
        else:
            assert isinstance(name, str), f"The first argument must a string: {name}"
            # if not extension, add ".py.md" as the extension
            if not os.path.splitext(name)[1]:
                name = f"{name}.py.md"

            with open(self.save_filepath(name), "w") as f:
                f.write(md_string)

    save_multiple_pkgs_code = multiple_pkgs_code  # backwards compatibility alias
